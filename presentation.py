from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]  # 0 - title and subtitle layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "TrekGuide Connect"
subtitle.text = "Your Ultimate Trekking Companion"

# Slide 2: Introduction
slide_layout = prs.slide_layouts[1]  # 1 - title and content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = ("- Revolutionizing the trekking experience.\n"
                "- Connects adventure enthusiasts with experienced trekking guides.\n"
                "- Offers a commission-based business model for entrepreneurs.")

# Slide 3: Comprehensive Guide Profiles
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Comprehensive Guide Profiles"
content.text = ("- Detailed profiles of registered trekking guides.\n"
                "- Information includes:\n"
                "  - Experience level\n"
                "  - Languages spoken\n"
                "  - Trekking routes expertise\n"
                "  - Customer reviews")

# Slide 4: Customized Matching Algorithm
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Customized Matching Algorithm"
content.text = ("- Advanced algorithm matches guides with user preferences and trekking requirements.\n"
                "- Suggests the most suitable guides.\n"
                "- Entrepreneurs earn commissions on successful bookings.")

# Slide 5: Interactive User Interface
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Interactive User Interface"
content.text = ("- User-friendly interface for browsing and comparing guides.\n"
                "- Filters and search options based on:\n"
                "  - Location\n"
                "  - Language\n"
                "  - Availability\n"
                "  - Trekking route")

# Slide 6: Secure Booking System
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Secure Booking System"
content.text = ("- Facilitates secure bookings through the platform.\n"
                "- Protects personal information and payment details.\n"
                "- Entrepreneurs earn commissions on each transaction.")

# Slide 7: Real-time Communication
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Real-time Communication"
content.text = ("- Integrated messaging system for seamless coordination.\n"
                "- Ensures all parties are well-informed throughout the trekking process.")

# Slide 8: Feedback and Reviews
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Feedback and Reviews"
content.text = ("- Clients can provide feedback and rate guides after treks.\n"
                "- Helps future trekkers make informed decisions.\n"
                "- Entrepreneurs use feedback to enhance service quality and attract clients.")

# Slide 9: Commission-Based Business Model
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Commission-Based Business Model"
content.text = ("- Entrepreneurs earn commissions on successful bookings.\n"
                "- Register as affiliates and receive a percentage of the total booking value.\n"
                "- Incentivizes promotion of the platform and facilitation of bookings.")

# Slide 10: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = ("- Simplifies trekking experience for adventure enthusiasts.\n"
                "- Provides a lucrative business opportunity for entrepreneurs.\n"
                "- Leverage the platform to capitalize on the growing demand for guided trekking tours.\n"
                "- Join TrekGuide Connect and embark on a journey to success in the booming tourism industry.")

# Save the presentation
file_path = "./TrekGuide_Connect_Presentation.pptx"
prs.save(file_path)
file_path
